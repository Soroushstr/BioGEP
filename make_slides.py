from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── colour palette ──────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x1F, 0x35, 0x64)   # slide titles / header rows
MID_BLUE    = RGBColor(0x2E, 0x75, 0xB6)   # section accents / column headers
LIGHT_BLUE  = RGBColor(0xBD, 0xD7, 0xEE)   # alternating table rows
GREEN       = RGBColor(0x37, 0x86, 0x44)   # positive highlight
ORANGE      = RGBColor(0xC5, 0x5A, 0x11)   # warning / problem
GREY_BG     = RGBColor(0xF2, 0xF2, 0xF2)   # slide background
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x00, 0x00, 0x00)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank


# ── helpers ─────────────────────────────────────────────────────────────────

def add_slide():
    slide = prs.slides.add_slide(BLANK)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = GREY_BG
    return slide


def txb(slide, text, left, top, width, height,
        size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT,
        italic=False, wrap=True):
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def title_bar(slide, title, subtitle=None):
    # full-width dark bar at top
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0), Inches(13.33), Inches(1.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = WHITE
    # left margin
    from pptx.util import Emu
    tf.margin_left = Emu(457200)   # 0.5 inch

    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(14)
        r2.font.bold = False
        r2.font.color.rgb = LIGHT_BLUE


def bullet_box(slide, items, left, top, width, height,
               size=16, indent=False):
    """items: list of str.  Lines starting with '•' get a bullet colour."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        if indent:
            p.level = 1
        bullet = item.startswith("•")
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.bold = item.startswith("▶")
        run.font.color.rgb = MID_BLUE if bullet else BLACK


def add_rect(slide, left, top, w, h, fill_rgb, line=False):
    r = slide.shapes.add_shape(1, Inches(left), Inches(top),
                               Inches(w), Inches(h))
    r.fill.solid()
    r.fill.fore_color.rgb = fill_rgb
    if line:
        r.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        r.line.width = Pt(0.5)
    else:
        r.line.fill.background()
    return r


def table(slide, headers, rows, left, top, width, height,
          col_widths=None):
    """Simple table drawn with rectangles and text boxes."""
    n_cols = len(headers)
    n_rows = len(rows)
    total_rows = n_rows + 1  # +1 header
    row_h = height / total_rows
    if col_widths is None:
        col_widths = [width / n_cols] * n_cols

    x = left
    # header row
    for ci, (hdr, cw) in enumerate(zip(headers, col_widths)):
        add_rect(slide, x, top, cw, row_h, MID_BLUE)
        txb(slide, hdr, x + 0.04, top + 0.02, cw - 0.08, row_h - 0.04,
            size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        x += cw

    for ri, row in enumerate(rows):
        y = top + row_h * (ri + 1)
        fill = LIGHT_BLUE if ri % 2 == 0 else WHITE
        x = left
        for ci, (cell, cw) in enumerate(zip(row, col_widths)):
            add_rect(slide, x, y, cw, row_h, fill, line=True)
            color = BLACK
            # highlight green / orange based on content
            try:
                val = float(str(cell).replace('**',''))
                if val >= 0.20:
                    color = GREEN
                elif val <= 0.05:
                    color = ORANGE
            except ValueError:
                pass
            if str(cell).startswith('**'):
                color = DARK_BLUE
            txb(slide, str(cell).replace('**',''), x + 0.04, y + 0.02,
                cw - 0.08, row_h - 0.04, size=12, color=color,
                align=PP_ALIGN.CENTER)
            x += cw


def divider(slide, y=1.15):
    """Thin accent line below title bar."""
    line = slide.shapes.add_shape(1, Inches(0), Inches(y),
                                  Inches(13.33), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = MID_BLUE
    line.line.fill.background()


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════════════════
slide = add_slide()
# big coloured background block
add_rect(slide, 0, 0, 13.33, 7.5, DARK_BLUE)
add_rect(slide, 0, 4.8, 13.33, 2.7, MID_BLUE)

txb(slide, "Cross-Species Gene Essentiality Prediction",
    0.6, 1.2, 12.0, 1.4, size=40, bold=True, color=WHITE)
txb(slide, "Using k-mer Graph Neural Networks with Domain Adaptation",
    0.6, 2.7, 12.0, 0.8, size=22, color=LIGHT_BLUE)
txb(slide, "Full Experiment History & Results Summary",
    0.6, 3.4, 12.0, 0.7, size=18, italic=True, color=LIGHT_BLUE)
txb(slide, "8 Species  ·  LOO Evaluation  ·  v5 → v10",
    0.6, 5.1, 10.0, 0.6, size=16, color=WHITE)
txb(slide, "April 2026",
    0.6, 5.8, 4.0, 0.5, size=14, color=LIGHT_BLUE)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — OVERVIEW / ROADMAP
# ════════════════════════════════════════════════════════════════════════════
slide = add_slide()
title_bar(slide, "Project Overview", "What are we trying to do?")
divider(slide)

txb(slide, "Goal", 0.4, 1.3, 2.5, 0.4, size=15, bold=True, color=MID_BLUE)
txb(slide,
    "Build a GNN that predicts gene essentiality and transfers across species — "
    "train on 7 species, predict the 8th (Leave-One-Out).",
    0.4, 1.7, 12.5, 0.7, size=15)

txb(slide, "8 Species", 0.4, 2.55, 2.5, 0.4, size=15, bold=True, color=MID_BLUE)
species_line = ("Arabidopsis  ·  Bacillus  ·  C. elegans  ·  Maripaludis  ·  "
                "Melanogaster  ·  Musculus  ·  Saccharomyces  ·  Sapiens")
txb(slide, species_line, 0.4, 2.95, 12.5, 0.5, size=14, italic=True)

txb(slide, "Architecture (fixed throughout)", 0.4, 3.6, 5.0, 0.4,
    size=15, bold=True, color=MID_BLUE)
arch = [
    "• GIN × 4 layers  |  hidden_dim = 256  |  k = 4 k-mer node features",
    "• Gradient Reversal Layer (GRL) for domain adaptation",
    "• Focal loss (γ=2)  |  Bayesian prior correction  |  lr = 3×10⁻⁴",
    "• Within-species sanity check (elegans → elegans): MCC = 0.87 ✓",
]
bullet_box(slide, arch, 0.4, 4.05, 12.5, 2.0, size=14)

# experiment roadmap boxes
steps = ["Step 1\nBaseline", "Step 2\nv7/v8\nProxy-val",
         "Step 3\nDiagnosis\nProxy-val bug",
         "Step 4\nv9\nLOO-val fix",
         "Step 5\nDiagnosis\nEpoch-1 bug",
         "Step 6\nv10\nConstant GRL",
         "Step 7\nConclusion"]
colors = [DARK_BLUE, MID_BLUE, ORANGE, GREEN, ORANGE, MID_BLUE, DARK_BLUE]
box_w = 1.7
gap = 0.15
start_x = 0.3
y_box = 5.9
for i, (s, c) in enumerate(zip(steps, colors)):
    x = start_x + i * (box_w + gap)
    add_rect(slide, x, y_box, box_w, 1.3, c)
    txb(slide, s, x + 0.05, y_box + 0.1, box_w - 0.1, 1.1,
        size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — STEP 1: INITIAL BASELINE
# ════════════════════════════════════════════════════════════════════════════
slide = add_slide()
title_bar(slide, "Step 1 — Initial Baseline (cfix_5 / cfix_6)",
          "First cross-species test")
divider(slide)

col1_x, col2_x = 0.4, 6.8
col_w = 6.0

# LEFT — what / expected
txb(slide, "Objective", col1_x, 1.3, col_w, 0.35, size=14, bold=True, color=MID_BLUE)
bullet_box(slide, [
    "• Build k-mer co-occurrence GNN (k=4)",
    "• Add biological node features (GC content, length…)",
    "• Add Gradient Reversal Layer (GRL) for domain adaptation",
    "• Test cross-species transfer for the first time",
], col1_x, 1.65, col_w, 1.6, size=13)

txb(slide, "Expected", col1_x, 3.35, col_w, 0.35, size=14, bold=True, color=MID_BLUE)
bullet_box(slide, [
    "• Some cross-species transfer (MCC > 0)",
    "• Within-species performance as a sanity check",
], col1_x, 3.7, col_w, 0.8, size=13)

# RIGHT — results table
txb(slide, "Results", col2_x, 1.3, col_w, 0.35, size=14, bold=True, color=MID_BLUE)
table(slide,
      ["Setting", "MCC", "AUC"],
      [
          ["Within-species: elegans → elegans", "0.870 ✓", "0.963"],
          ["Cross-species: all → saccharomyces (v5)", "0.071", "0.551"],
          ["Cross-species: all → saccharomyces (v6)", "0.128", "0.593"],
      ],
      col2_x, 1.7, col_w, 1.7,
      col_widths=[3.0, 1.5, 1.5])

txb(slide, "Key Takeaway", col2_x, 3.5, col_w, 0.35, size=14, bold=True, color=ORANGE)
bullet_box(slide, [
    "• Within-species MCC = 0.87 — architecture can learn essentiality ✓",
    "• Cross-species MCC ≈ 0.07–0.13 — transfer is hard but non-zero",
    "• Large gap between within-species and cross-species performance",
    "  confirms the domain shift challenge",
], col2_x, 3.85, col_w, 1.8, size=13)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — STEP 2: v7 / v8 PROXY-VAL
# ════════════════════════════════════════════════════════════════════════════
slide = add_slide()
title_bar(slide, "Step 2 — Full LOO Evaluation: v7 & v8",
          "Proxy-val early stopping — training on 7 species, testing on 1")
divider(slide)

txb(slide, "Setup", 0.4, 1.25, 4.0, 0.35, size=14, bold=True, color=MID_BLUE)
bullet_box(slide, [
    "• Full 8-species LOO evaluation (each species held out once)",
    "• Validation set drawn from training species (proxy-val)",
    "• MCC-based early stopping  |  up to 200 epochs",
], 0.4, 1.6, 5.5, 1.0, size=13)

# v7 table
txb(slide, "v7 Results (proxy-val, MCC stopping)  —  Mean MCC = 0.119",
    0.4, 2.7, 6.0, 0.38, size=13, bold=True, color=DARK_BLUE)
table(slide,
      ["Species", "Test MCC", "Test AUC", "Best Epoch"],
      [
          ["arabidopsis", "0.191", "0.595", "190"],
          ["bacillus",    "0.109", "0.599", "189"],
          ["elegans",     "0.098", "0.651", "199"],
          ["maripaludis", "0.088", "0.547", "183"],
          ["melanogaster","0.098", "0.627", "188"],
          ["musculus",    "0.143", "0.618", "24"],
          ["saccharomyces","0.058","0.540", "188"],
          ["sapiens",     "0.166", "0.668", "186"],
      ],
      0.4, 3.1, 5.9, 3.8,
      col_widths=[1.9, 1.2, 1.2, 1.6])

# v8 table
txb(slide, "v8 Results (proxy-val)  —  Mean MCC = 0.116",
    6.5, 2.7, 6.4, 0.38, size=13, bold=True, color=DARK_BLUE)
table(slide,
      ["Species", "Test MCC", "Test AUC", "Best Epoch"],
      [
          ["arabidopsis", "0.137", "0.562", "177"],
          ["bacillus",    "0.125", "0.624", "196"],
          ["elegans",     "0.085", "0.652", "189"],
          ["maripaludis", "0.095", "0.554", "177"],
          ["melanogaster","0.082", "0.606", "200"],
          ["musculus",    "0.144", "0.617", "27"],
          ["saccharomyces","0.085","0.557", "197"],
          ["sapiens",     "0.180", "0.675", "199"],
      ],
      6.5, 3.1, 5.9, 3.8,
      col_widths=[1.9, 1.2, 1.2, 1.6])

# warning box
add_rect(slide, 0.4, 6.6, 12.5, 0.65, ORANGE)
txb(slide,
    "⚠  Val MCC at best epochs ≈ 0.65–0.70  vs  Test MCC ≈ 0.08–0.19  →  3–12× gap!  "
    "Most best epochs are ~180–200, meaning the model trains to near-completion on training-species val.",
    0.55, 6.65, 12.2, 0.55, size=12, bold=True, color=WHITE)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — STEP 3: DIAGNOSIS — PROXY-VAL BUG
# ════════════════════════════════════════════════════════════════════════════
slide = add_slide()
title_bar(slide, "Step 3 — Diagnosing the Proxy-Val Problem",
          "Why is val MCC 3–12× higher than test MCC?")
divider(slide)

# diagram: two boxes connected by arrow
add_rect(slide, 0.5, 1.4, 4.5, 1.5, MID_BLUE)
txb(slide, "Validation Set\n(proxy-val)",
    0.55, 1.45, 4.4, 0.5, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(slide,
    "Drawn from the same 7\ntraining species",
    0.55, 1.95, 4.4, 0.9, size=13, color=WHITE, align=PP_ALIGN.CENTER)

txb(slide, "→", 5.1, 1.85, 0.6, 0.8, size=30, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

add_rect(slide, 5.8, 1.4, 4.5, 1.5, ORANGE)
txb(slide, "Early Stopping\n(MCC on proxy-val)",
    5.85, 1.45, 4.4, 0.5, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(slide,
    "Stops when training-species\nperformance peaks (~ep 185)",
    5.85, 1.95, 4.4, 0.9, size=13, color=WHITE, align=PP_ALIGN.CENTER)

txb(slide, "→", 10.4, 1.85, 0.6, 0.8, size=30, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

add_rect(slide, 11.1, 1.4, 1.9, 1.5, RGBColor(0xC0, 0x00, 0x00))
txb(slide, "Model that\noverfits\ntraining species",
    11.15, 1.45, 1.8, 1.4, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

bullet_box(slide, [
    "▶ What actually happens:",
    "• The validation set is from the 7 training species → val MCC reflects training-species performance.",
    "• Model keeps improving on training species → val MCC keeps rising → training runs to epoch ~185.",
    "• Meanwhile, generalisation to the held-out species peaks earlier and then degrades.",
    "• Result: we always pick the WORST model for cross-species transfer.",
    "",
    "▶ The fix:  use data from the held-out species itself as validation during training.",
    "• Reserve 10% or 20% of the test species before training begins.",
    "• Val MCC will now track cross-species generalisation — stopping will be meaningful.",
], 0.4, 3.2, 12.5, 3.8, size=14)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — STEP 4: v9 LOO-VAL FIX
# ════════════════════════════════════════════════════════════════════════════
slide = add_slide()
title_bar(slide, "Step 4 — v9: LOO-Val Fix",
          "Validation from the held-out species itself (10% and 20%)")
divider(slide)

txb(slide, "Change made", 0.4, 1.25, 5.0, 0.35, size=14, bold=True, color=MID_BLUE)
bullet_box(slide, [
    "• Reserve 10% (lv10) or 20% (lv20) of test species as stratified val set",
    "• Early stopping and model selection now based on cross-species generalisation signal",
    "• MCC-based stopping  |  patience = 20",
], 0.4, 1.6, 6.0, 1.0, size=13)

# lv20 table
txb(slide, "v9 lv20 (20% val)  —  Mean MCC = 0.174  (+46% vs v7)",
    0.4, 2.7, 6.0, 0.38, size=13, bold=True, color=GREEN)
table(slide,
      ["Species", "Test MCC", "Test AUC", "Best Epoch"],
      [
          ["arabidopsis",  "0.083", "0.535", "46"],
          ["bacillus",     "0.224", "0.670", "1"],
          ["elegans",      "0.196", "0.735", "3"],
          ["maripaludis",  "0.252", "0.660", "1"],
          ["melanogaster", "0.133", "0.669", "4"],
          ["musculus",     "0.166", "0.632", "1"],
          ["saccharomyces","0.102", "0.561", "8"],
          ["sapiens",      "0.234", "0.721", "2"],
      ],
      0.4, 3.1, 5.9, 3.8,
      col_widths=[1.9, 1.2, 1.2, 1.6])

# lv10 table
txb(slide, "v9 lv10 (10% val)  —  Mean MCC = 0.169  (+42% vs v7)",
    6.5, 2.7, 6.4, 0.38, size=13, bold=True, color=GREEN)
table(slide,
      ["Species", "Test MCC", "Test AUC", "Best Epoch"],
      [
          ["arabidopsis",  "0.101", "0.562", "68"],
          ["bacillus",     "0.189", "0.657", "8"],
          ["elegans",      "0.217", "0.747", "6"],
          ["maripaludis",  "0.248", "0.656", "1"],
          ["melanogaster", "0.132", "0.659", "11"],
          ["musculus",     "0.146", "0.616", "38"],
          ["saccharomyces","0.092", "0.562", "40"],
          ["sapiens",      "0.229", "0.723", "1"],
      ],
      6.5, 3.1, 5.9, 3.8,
      col_widths=[1.9, 1.2, 1.2, 1.6])

# result highlight + new problem
add_rect(slide, 0.4, 6.6, 5.9, 0.65, GREEN)
txb(slide, "✓  Val-test gap closed  |  Mean MCC: 0.119 → 0.174  (+46%)",
    0.55, 6.65, 5.7, 0.55, size=13, bold=True, color=WHITE)

add_rect(slide, 6.5, 6.6, 6.5, 0.65, ORANGE)
txb(slide, "⚠  NEW PROBLEM: Best epoch = 1–8 for most species — model generalises best before training!",
    6.6, 6.65, 6.3, 0.55, size=12, bold=True, color=WHITE)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — STEP 5: DIAGNOSIS — GRL SCHEDULE BUG
# ════════════════════════════════════════════════════════════════════════════
slide = add_slide()
title_bar(slide, "Step 5 — Diagnosing the 'Best Epoch = 1' Problem",
          "Why does the model generalise best before training?")
divider(slide)

txb(slide, "Observation", 0.4, 1.3, 12.5, 0.35, size=14, bold=True, color=MID_BLUE)
txb(slide,
    "For 5–7 out of 8 species, the held-out species val AUC/MCC peaks at epoch 1 "
    "and then monotonically decreases throughout training. The model is not learning "
    "generalisable features — something is actively hurting cross-species transfer.",
    0.4, 1.65, 12.5, 0.7, size=13)

txb(slide, "Root Cause — GRL Alpha Schedule", 0.4, 2.5, 8.0, 0.38, size=14, bold=True, color=ORANGE)
add_rect(slide, 0.4, 2.9, 8.5, 1.2, RGBColor(0xFF, 0xF2, 0xCC))
txb(slide,
    'grl_alpha  =  2 / (1 + exp(−10 × epoch / max_epochs)) − 1\n\n'
    'Epoch 1 of 200:  grl_alpha ≈ 0.000    |    Epoch 10 of 200:  grl_alpha ≈ 0.222',
    0.55, 2.95, 8.2, 1.1, size=13, color=RGBColor(0x7F, 0x3F, 0x00))

bullet_box(slide, [
    "▶ What this means:",
    "• The adversarial (GRL) constraint is essentially zero for the first ~20 epochs.",
    "• During those epochs the model freely learns species-specific features with no penalty.",
    "• The more training occurs, the more species-specific the representations become.",
    "• Val MCC on the held-out species falls after epoch 1 because the features stop being universal.",
    "",
    "▶ The fix for v10:",
    "• Set  grl_alpha = 1.0  (constant) — adversarial constraint active from epoch 1.",
    "• Test three lambda (adversarial loss weight) values: 0.5, 1.0, 2.0.",
    "• Switch to AUC-based stopping (more stable than MCC on small val sets)  |  patience = 20.",
], 0.4, 4.2, 12.5, 3.0, size=13)

# decision rule box
add_rect(slide, 0.4, 6.55, 12.5, 0.75, DARK_BLUE)
txb(slide,
    "Decision rule:  Best epochs shift to 20+ AND MCC improves → GRL works.  "
    "Best epochs still 1–5 at λ=2.0 → GRL is structurally insufficient → Priority 2 (phylogeny-aware training).",
    0.55, 6.6, 12.2, 0.65, size=12, bold=True, color=WHITE)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — STEP 6: v10 RESULTS
# ════════════════════════════════════════════════════════════════════════════
slide = add_slide()
title_bar(slide, "Step 6 — v10: Constant GRL  +  AUC Stopping",
          "λ = 0.5 (a05), 1.0 (a10), 2.0 (a20)  —  24 SLURM jobs")
divider(slide)

# combined table all 3 lambdas
txb(slide, "Results across all three lambda variants (best epoch shown):",
    0.4, 1.25, 12.5, 0.38, size=14, bold=True, color=DARK_BLUE)

table(slide,
      ["Species",
       "v9 MCC", "v9 Ep",
       "a05 MCC", "a05 Ep",
       "a10 MCC", "a10 Ep",
       "a20 MCC", "a20 Ep"],
      [
          ["arabidopsis",  "0.083","46",  "0.089","18", "0.118","1",  "0.057","43"],
          ["bacillus",     "0.224","1",   "0.213","1",  "0.179","1",  "0.177","1"],
          ["elegans",      "0.196","3",   "0.191","6",  "0.193","6",  "0.189","3"],
          ["maripaludis",  "0.252","1",   "0.243","1",  "0.233","1",  "0.259","1"],
          ["melanogaster", "0.133","4",   "0.127","3",  "0.155","5",  "0.144","5"],
          ["musculus",     "0.166","1",   "0.169","1",  "0.166","2",  "0.178","1"],
          ["saccharomyces","0.102","8",   "0.123","1",  "0.125","1",  "0.126","18"],
          ["sapiens",      "0.234","2",   "0.245","1",  "0.245","1",  "0.243","1"],
          ["MEAN",         "0.174","",    "0.175","",   "0.164","",   "0.172",""],
      ],
      0.4, 1.65, 12.5, 4.55,
      col_widths=[1.55, 0.85, 0.7, 0.85, 0.7, 0.85, 0.7, 0.85, 0.7])

# verdict
add_rect(slide, 0.4, 6.35, 5.9, 0.9, ORANGE)
txb(slide,
    "⚠  5/8 species still stop at epoch 1 even at λ=2.0\n"
    "Best epoch pattern essentially unchanged from v9",
    0.55, 6.38, 5.7, 0.84, size=12, bold=True, color=WHITE)

add_rect(slide, 6.5, 6.35, 6.5, 0.9, RGBColor(0xC0, 0x00, 0x00))
txb(slide,
    "✗  Constant GRL did NOT fix the problem.\n"
    "MCC improvement is negligible (best: +0.1% at λ=0.5).\n"
    "Decision rule → GRL is insufficient → Priority 2.",
    6.6, 6.38, 6.3, 0.84, size=12, bold=True, color=WHITE)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — OVERALL COMPARISON
# ════════════════════════════════════════════════════════════════════════════
slide = add_slide()
title_bar(slide, "Overall Results — All Versions Compared",
          "Mean Test MCC across 8 species (LOO)")
divider(slide)

table(slide,
      ["Version", "Key Change", "Mean MCC", "Δ vs v7"],
      [
          ["v7", "Full LOO eval, proxy-val MCC stopping", "0.119", "baseline"],
          ["v8", "Same setup, more epochs (200)", "0.116", "≈ 0"],
          ["v9 lv10", "LOO-val 10% held-out, MCC stopping", "0.169", "+42% ✓"],
          ["v9 lv20", "LOO-val 20% held-out, MCC stopping", "0.174", "+46% ✓"],
          ["v10 a05", "Constant GRL λ=0.5, AUC stopping", "0.175", "+47% (≈v9)"],
          ["v10 a10", "Constant GRL λ=1.0, AUC stopping", "0.164", "+38%"],
          ["v10 a20", "Constant GRL λ=2.0, AUC stopping", "0.172", "+44%"],
      ],
      0.4, 1.3, 12.5, 3.6,
      col_widths=[1.5, 5.5, 1.8, 3.7])

txb(slide, "Key lessons", 0.4, 5.1, 12.5, 0.38, size=14, bold=True, color=DARK_BLUE)
bullet_box(slide, [
    "• The single biggest gain came from fixing the validation strategy: proxy-val → LOO-val  (+46%)",
    "• GRL in any form (scheduled or constant, any λ) did not improve cross-species MCC meaningfully",
    "• The 'best epoch = 1' symptom persists — training actively hurts cross-species transfer",
    "• Root cause is structural: phylogenetic distance between training and test species is too large",
], 0.4, 5.5, 12.5, 1.7, size=14)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CONCLUSIONS & NEXT STEPS
# ════════════════════════════════════════════════════════════════════════════
slide = add_slide()
title_bar(slide, "Conclusions & Next Steps",
          "Where we are and where to go")
divider(slide)

txb(slide, "What worked", 0.4, 1.3, 4.0, 0.38, size=14, bold=True, color=GREEN)
bullet_box(slide, [
    "✓  Architecture learns essentiality (within-species MCC = 0.87)",
    "✓  Fixing val strategy (LOO-val 20%) gave +46% MCC improvement",
    "✓  Val-test gap closed — stopping is now meaningful",
    "✓  Best model: v9 lv20 / v10 a05  |  Mean test MCC ≈ 0.174–0.175",
], 0.4, 1.7, 5.8, 1.8, size=13)

txb(slide, "What didn't work", 0.4, 3.6, 5.0, 0.38, size=14, bold=True, color=ORANGE)
bullet_box(slide, [
    "✗  Proxy-val stopping  (3–12× val-test gap)",
    "✗  Scheduled GRL alpha  (near-zero for first 20 epochs)",
    "✗  Constant GRL at any lambda  (best epochs still = 1)",
    "✗  More training epochs alone  (v8 ≈ v7)",
], 0.4, 4.0, 5.8, 1.8, size=13)

# Priority 2 box
add_rect(slide, 6.6, 1.25, 6.3, 5.1, DARK_BLUE)
txb(slide, "Priority 2 — Phylogeny-Aware Training",
    6.75, 1.35, 6.1, 0.5, size=15, bold=True, color=WHITE)
txb(slide,
    "Core insight: training on 8 phylogenetically diverse species forces "
    "the model to learn features that are universal across bacteria, archaea, "
    "plants, fungi, worms, flies, and mammals — but no such universal signal exists "
    "at the k-mer level.",
    6.75, 1.9, 6.1, 1.1, size=12, color=LIGHT_BLUE)
txb(slide, "Proposed fix:", 6.75, 3.1, 6.1, 0.35, size=13, bold=True, color=WHITE)
bullet_box(slide, [
    "• For each held-out species, select only",
    "  phylogenetically close training species.",
    "• E.g., hold out sapiens → train on",
    "  {musculus, elegans, melanogaster}",
    "• E.g., hold out bacillus → train on",
    "  {maripaludis, saccharomyces}",
    "• Reduces noise from unrelated species,",
    "  sharpens cross-species signal.",
], 6.75, 3.5, 6.1, 2.5, size=12)

add_rect(slide, 0.4, 6.55, 12.5, 0.75, MID_BLUE)
txb(slide,
    "Current best MCC = 0.174  |  Target: meaningful improvement via phylogeny-aware training  |  "
    "GRL approaches exhausted",
    0.55, 6.6, 12.2, 0.65, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════════════════
out = "/gpfs/data/fs72129/imansamiei/soroush/GEPNew/newsrc/newsrc/src/GEP_Experiment_Summary.pptx"
prs.save(out)
print(f"Saved: {out}")
